import os
from pptx import Presentation

def search_pptx_files(directory, search_text):
    found_files = set() # use set to ensure uniqueness
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith('.pptx'):
                pptx_file_path = os.path.join(root, file)
                try:
                    presentation = Presentation(pptx_file_path)
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                if search_text in shape.text:
                                    found_files.add(pptx_file_path)
                                    break
                    print('.', end='', flush=True) # print a dot after each file has been searched
                except Exception as e:
                    print(f"Error opening file {pptx_file_path}: {str(e)}")
    return found_files

# Get user input
search_text = input("Enter the text you want to search for: ")
directory = input("Enter the directory to search in: ")

# Search for pptx files containing the text
found_files = search_pptx_files(directory, search_text)

# Print found files
print("\nFound the text in following PowerPoint files:")
for file in found_files:
    print(file)

